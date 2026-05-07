"""Component-based PPTX export — crop each posterbox region from the
compiled PDF and embed each crop as an independent shape in PPTX.

This produces pixel-perfect fidelity (from PDF) while keeping each card
movable / resizable in PowerPoint.
"""

from __future__ import annotations

import os
import shutil
import tempfile

import fitz
from pptx import Presentation
from pptx.util import Mm
from pptx.dml.color import RGBColor


PDF_PATH = "main.pdf"
OUT_PATH = "poster_components.pptx"


def main():
    doc = fitz.open(PDF_PATH)
    page = doc[0]
    pw, ph = page.rect.width, page.rect.height

    # A0 landscape mm
    W_mm, H_mm = 1189, 841

    def pts_to_mm(x, y):
        return x / pw * W_mm, y / ph * H_mm

    COLS = 4
    ROWS = 20
    row_h = ph / ROWS
    col_w = pw / COLS

    # Match main.tex `between=rowN and rowM` boundaries
    regions = {
        "title":         (0, 0,  4, 4),
        "stats":         (0, 4,  4, 6),
        "background":    (0, 6,  1, 11),
        "contributions": (0, 11, 1, 16),
        "refs":          (0, 16, 1, 20),
        "method":        (1, 6,  1, 14),
        "mechanism":     (1, 14, 1, 20),
        "results1":      (2, 6,  1, 14),
        "phase":         (2, 14, 1, 20),
        "waterbirds":    (3, 6,  1, 14),
        "takeaways":     (3, 14, 1, 20),
    }

    prs = Presentation()
    prs.slide_width = Mm(W_mm)
    prs.slide_height = Mm(H_mm)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)  # ICLR light green

    tmpdir = tempfile.mkdtemp()
    mat = fitz.Matrix(300 / 72, 300 / 72)

    for name, (col, r0, span, r1) in regions.items():
        clip = fitz.Rect(col * col_w, r0 * row_h,
                         (col + span) * col_w, r1 * row_h)
        pix = page.get_pixmap(matrix=mat, clip=clip)
        img_path = os.path.join(tmpdir, f"{name}.png")
        pix.save(img_path)

        left, top = pts_to_mm(clip.x0, clip.y0)
        right, bottom = pts_to_mm(clip.x1, clip.y1)

        slide.shapes.add_picture(img_path, Mm(left), Mm(top),
                                 Mm(right - left), Mm(bottom - top))

    prs.save(OUT_PATH)
    doc.close()
    shutil.rmtree(tmpdir)
    print(f"[pptx] wrote {OUT_PATH} with {len(regions)} component shapes")


if __name__ == "__main__":
    main()
